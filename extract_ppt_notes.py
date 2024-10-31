from pptx import Presentation
import sys

def extract_notes(ppt_path, output_path):
    # 打开PPT文件
    prs = Presentation(ppt_path)
    
    # 打开输出文件
    with open(output_path, 'w', encoding='utf-8') as f:
        # 遍历所有幻灯片
        for slide_number, slide in enumerate(prs.slides, 1):
            # 获取备注
            notes = slide.notes_slide
            if notes and notes.notes_text_frame.text.strip():
                # 写入幻灯片编号和对应的备注
                #f.write(f"=== 第{slide_number}张幻灯片的备注 ===\n")
                f.write(notes.notes_text_frame.text.strip())
                f.write('\n\n')

def main():
    if len(sys.argv) != 3:
        print("使用方法: python extract_notes.py <PPT文件路径> <输出文件路径>")
        sys.exit(1)
    
    ppt_path = sys.argv[1]
    output_path = sys.argv[2]
    
    try:
        extract_notes(ppt_path, output_path)
        print(f"备注已成功导出到: {output_path}")
    except Exception as e:
        print(f"发生错误: {str(e)}")

if __name__ == "__main__":
    main()
